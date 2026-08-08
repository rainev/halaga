#!/usr/bin/env python3
"""Build FinSight-Business-Summary.pptx — a concise business & unit-economics deck.
Run: /opt/anaconda3/bin/python3 build_summary_deck.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

DARK = RGBColor(0x0B, 0x3D, 0x2E)
MID = RGBColor(0x10, 0x55, 0x3F)
ACCENT = RGBColor(0x1F, 0x7A, 0x58)
MINT = RGBColor(0x7C, 0xF0, 0xC0)
CREAM = RGBColor(0xEA, 0xFF, 0xF6)
INK = RGBColor(0x1A, 0x22, 0x30)
GREY = RGBColor(0x5B, 0x6B, 0x63)
LIGHT = RGBColor(0xF1, 0xF7, 0xF3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _tf(tb):
    tf = tb.text_frame
    tf.word_wrap = True
    return tf


def box(slide, x, y, w, h):
    return slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))


def fill(slide, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def set_text(tf, text, size, color, bold=False, align=PP_ALIGN.LEFT, italic=False):
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color; f.name = "Helvetica Neue"
    return p


def bullets(slide, x, y, w, h, items, size=17, color=INK, gap=6):
    tb = box(slide, x, y, w, h); tf = _tf(tb)
    for i, (txt, bold) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        r = p.add_run(); r.text = "•  " + txt
        f = r.font; f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = "Helvetica Neue"
    return tb


def header(slide, kicker, title):
    fill(slide, 0, 0, 13.333/72*72/72*0, 0, DARK)  # noop guard
    fill(slide, 0.0, 0.0, 13.333, 1.55, DARK)
    fill(slide, 0.7, 1.28, 1.4, 0.06, MINT)
    tb = box(slide, 0.7, 0.28, 12, 0.4); set_text(_tf(tb), kicker.upper(), 12, MINT, bold=True)
    tb2 = box(slide, 0.7, 0.55, 12, 0.8); set_text(_tf(tb2), title, 30, WHITE, bold=True)


def bignum(slide, x, y, num, label, w=4.0):
    tb = box(slide, x, y, w, 1.0); set_text(_tf(tb), num, 40, ACCENT, bold=True)
    tb2 = box(slide, x, y + 0.85, w, 0.7); set_text(_tf(tb2), label, 14, GREY)


def table(slide, x, y, w, rows, colw, header_row=True, fs=13, rowh=0.42):
    ncol = len(rows[0])
    tbl_h = rowh * len(rows)
    gt = slide.shapes.add_table(len(rows), ncol, Inches(x), Inches(y), Inches(w), Inches(tbl_h)).table
    for ci, cw in enumerate(colw):
        gt.columns[ci].width = Inches(cw)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = gt.cell(ri, ci)
            cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; r = p.add_run(); r.text = str(val)
            f = r.font; f.size = Pt(fs); f.name = "Helvetica Neue"
            if header_row and ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = MID
                f.color.rgb = WHITE; f.bold = True
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
                f.color.rgb = INK
            if ci >= 1 and not (header_row and ri == 0):
                p.alignment = PP_ALIGN.RIGHT
    return gt


def footer(slide, n):
    tb = box(slide, 11.6, 7.05, 1.6, 0.3)
    set_text(_tf(tb), f"FinSight · {n}", 9, GREY, align=PP_ALIGN.RIGHT)


# ---------- Slide 1 — cover ----------
s = prs.slides.add_slide(BLANK)
fill(s, 0, 0, 13.333, 7.5, DARK)
fill(s, 0.9, 3.5, 2.2, 0.09, MINT)
tb = box(s, 0.9, 2.35, 11, 0.5); set_text(_tf(tb), "FINSIGHT", 16, MINT, bold=True)
tb = box(s, 0.85, 2.75, 11.6, 1.2); set_text(_tf(tb), "Business Summary", 52, WHITE, bold=True)
tb = box(s, 0.9, 3.75, 11, 0.8)
set_text(_tf(tb), "PSE valuation + portfolio-aware insights — unit economics, pricing & scaling", 18, CREAM)
tb = box(s, 0.9, 6.5, 11, 0.5); set_text(_tf(tb), "Draft · July 2026", 13, RGBColor(0x9F, 0xC8, 0xB8))

# ---------- Slide 2 — what it is ----------
s = prs.slides.add_slide(BLANK); header(s, "The product", "What FinSight is")
bullets(s, 0.7, 1.95, 12, 4.5, [
    ("An investment-awareness platform for Philippine Stock Exchange (PSE) investors — a live app.", True),
    ("Pillar 1 — Valuation workbench: value any PSE stock with 4 classic models (DCF, DDM, Graham, Multiples), PH-calibrated.", False),
    ("Pillar 2 — Portfolio-aware insights: “never be blindsided by news that touches what you own.”", False),
    ("The contract: awareness, not advice — we surface, connect, and cite the source; never “buy/sell.”", False),
    ("Insights are generated once per (article, company) and shared across all users — so AI cost is ~flat.", False),
], size=18, gap=12)
footer(s, "2")

# ---------- Slide 3 — market ----------
s = prs.slides.add_slide(BLANK); header(s, "The market", "A market in a land-grab moment")
bignum(s, 0.7, 2.0, "2.86M", "PSE accounts in 2024 — +50% in one year", w=3.6)
bignum(s, 4.6, 2.0, "86%", "of accounts are online / mobile", w=3.2)
bignum(s, 8.0, 2.0, "53%", "of adults even aware of investing", w=4.0)
bullets(s, 0.7, 3.9, 12, 2.6, [
    ("18–29 is the fastest-growing cohort (26.5%, up from 19.5%) — learns to invest on TikTok.", False),
    ("New, mobile-first, socially-taught and under-equipped — exactly FinSight's customer.", True),
    ("Source: PSE 2024 profile · DataReportal 2025 · BSP CFIS (see market-research.md).", False),
], size=16, gap=10)
footer(s, "3")

# ---------- Slide 4 — pricing / paywall ----------
s = prs.slides.add_slide(BLANK); header(s, "The model", "How we charge — 3 months free, then pay")
tb = box(s, 0.7, 1.75, 12, 0.4)
set_text(_tf(tb), "Two tiers · four billing terms · everyone pays after the free trial", 15, ACCENT, bold=True)
table(s, 0.7, 2.35, 5.9, [
    ["Standard — ₱99/mo", "Price", "/mo"],
    ["Monthly", "₱99", "₱99"],
    ["3 months", "₱279", "₱93"],
    ["6 months", "₱534", "₱89"],
    ["Annual (2 mo free)", "₱990", "₱82.50"],
], [3.0, 1.5, 1.4], fs=13)
table(s, 6.9, 2.35, 5.7, [
    ["Pro — ₱299/mo", "Price", "/mo"],
    ["Monthly", "₱299", "₱299"],
    ["3 months", "₱849", "₱283"],
    ["6 months", "₱1,614", "₱269"],
    ["Annual (2 mo free)", "₱2,990", "₱249"],
], [3.0, 1.4, 1.3], fs=13)
bullets(s, 0.7, 5.25, 12, 1.6, [
    ("Longer terms pull cash forward and cut churn; Pro lifts ARPU without raising the ₱99 entry point.", True),
    ("Pro adds: unlimited holdings, real-time alerts, sector/peer insights, history/export, priority support.", False),
], size=15, gap=9)
footer(s, "4")

# ---------- Slide 5 — unit economics ----------
s = prs.slides.add_slide(BLANK); header(s, "Unit economics", "Everyone pays ₱99 — the model that works")
bignum(s, 0.7, 2.0, "₱99", "revenue per user / mo (vs ~₱15 for freemium)", w=3.4)
bignum(s, 4.3, 2.0, "~20", "paying users to break even", w=3.2)
bignum(s, 7.6, 2.0, "~₱13", "tech cost per user / mo", w=4.2)
bullets(s, 0.7, 3.9, 12, 2.6, [
    ("Because everyone pays, revenue/user is ~6.6× a 5%-conversion freemium — lower price, but all pay.", True),
    ("Tech cost per user falls as you scale (₱13 → ₱7); AI is a rounding error (shared insights).", False),
    ("Cost is never the constraint — acquisition, trial conversion and churn are.", False),
], size=16, gap=10)
footer(s, "5")

# ---------- Slide 6 — headcount ladder ----------
s = prs.slides.add_slide(BLANK); header(s, "Staying lean", "Headcount grows far slower than users")
table(s, 0.7, 2.0, 7.4, [
    ["Users", "People", "Who"],
    ["10 – 2,000", "1", "just you (solo founder)"],
    ["3,000 – 5,000", "2", "you + support"],
    ["10,000", "3", "you + support + engineer"],
], [2.2, 1.6, 3.6], fs=15, rowh=0.6)
bullets(s, 8.5, 2.0, 4.4, 4.5, [
    ("Solo to ~2,000 users.", True),
    ("Hire only when revenue already pays for it.", False),
    ("Users grow 500×; people grow 1 → 3.", True),
    ("That gap is the whole business.", False),
], size=15, gap=10)
footer(s, "6")

# ---------- Slide 7 — profitability ----------
s = prs.slides.add_slide(BLANK); header(s, "Profitability", "Profit at ₱99, everyone pays")
table(s, 0.7, 1.95, 8.4, [
    ["Users", "Revenue / mo", "Burn / mo", "Profit / mo"],
    ["500", "₱49,500", "₱4,000", "+₱45,500"],
    ["1,000", "₱99,000", "₱5,000", "+₱94,000"],
    ["2,000", "₱198,000", "₱7,000", "+₱191,000"],
    ["5,000", "₱495,000", "₱43,000", "+₱452,000"],
    ["10,000", "₱990,000", "₱134,000", "+₱856,000"],
], [1.7, 2.3, 2.1, 2.3], fs=14, rowh=0.5)
tb = box(s, 9.4, 2.4, 3.4, 3)
set_text(_tf(tb), "Steady-state, everyone-pays view. At a realistic 50% trial→paid, halve it — still ~₱361k/mo profit at 10k users.", 14, GREY)
footer(s, "7")

# ---------- Slide 8 — annual cash ----------
s = prs.slides.add_slide(BLANK); header(s, "Cash flow", "Annual plans fund the whole business")
bignum(s, 0.7, 2.1, "₱9.9M", "cash collected upfront at 10,000 annual subs", w=5.0)
bignum(s, 7.2, 2.1, "~₱8.3M", "annual profit at 10k (₱9.9M − ₱1.6M burn)", w=5.5)
bullets(s, 0.7, 4.1, 12, 2.3, [
    ("Annual (₱990, “2 months free”) collects a full year at signup — self-funds growth & hiring.", True),
    ("Cuts churn to one renewal decision a year; boosts lifetime value.", False),
    ("Offer both monthly and annual at the paywall; 30–50% annual take-up is common.", False),
], size=16, gap=10)
footer(s, "8")

# ---------- Slide 9 — go to market ----------
s = prs.slides.add_slide(BLANK); header(s, "Go-to-market", "The only real constraint: getting users")
bullets(s, 0.7, 2.0, 12, 4.4, [
    ("Organic-first social: TikTok (lead), Facebook, Instagram — where PH's 18–29 investors already are.", True),
    ("Content pillars: “Halaga check” valuation series, “blindsided” awareness, and street/vox-pop interviews.", False),
    ("Message = awareness, not advice — the compliant, credible alternative to “hot tip” finfluencers.", False),
    ("Paid amplification later — boost the organic posts that already win. (See social-campaigns.md.)", False),
    ("Cost is trivial; the entire game is acquisition → trial conversion → retention.", True),
], size=17, gap=12)
footer(s, "9")

# ---------- Slide 10 — summary ----------
s = prs.slides.add_slide(BLANK)
fill(s, 0, 0, 13.333, 7.5, DARK)
fill(s, 0.9, 1.35, 2.0, 0.08, MINT)
tb = box(s, 0.9, 0.7, 11, 0.6); set_text(_tf(tb), "THE PUNCHLINE", 14, MINT, bold=True)
tb = box(s, 0.85, 1.55, 11.6, 0.9); set_text(_tf(tb), "Cheap to run, profitable early, cash-rich at scale", 30, WHITE, bold=True)
items = [
    "₱99/mo, everyone pays after a 3-month free trial — plus 3/6/12-month terms and a ₱299 Pro tier.",
    "Break even at ~20 paying users; solo and profitable to ~2,000 users.",
    "10,000 users → ~₱990k/mo revenue, ~₱856k/mo profit — with a 3-person team.",
    "Annual plans collect ~₱9.9M upfront at 10k subs — self-funds growth.",
    "Tech (cloud + news + AI) is a rounding error; people are the cost, hired behind revenue.",
    "The only constraint is acquisition — which the organic campaign plan targets directly.",
]
tb = box(s, 0.9, 2.6, 11.5, 4.3); tf = _tf(tb)
for i, t in enumerate(items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(11)
    r = p.add_run(); r.text = "→  " + t
    f = r.font; f.size = Pt(17); f.color.rgb = CREAM; f.name = "Helvetica Neue"

prs.save("FinSight-Business-Summary.pptx")
print("saved FinSight-Business-Summary.pptx")
